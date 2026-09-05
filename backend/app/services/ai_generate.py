"""Generate draf form (sections + soal + settings) via Gemini.

Alur: prompt + file referensi (teks) -> Gemini (JSON mode) -> draf yang
sudah disanitasi. Draf TIDAK langsung jadi form — creator mereview lalu
POST /api/ai/accept yang memvalidasi ulang memakai skema existing.
"""
import html
import io
import json
import logging
import re

import httpx

from app.config import GEMINI_API_KEY, GEMINI_FALLBACK_MODEL, GEMINI_MODEL
from app.schemas.question import QuestionCreate

logger = logging.getLogger("quizary.ai")

AI_DAILY_LIMIT = 3

ALLOWED_REF_EXT = {".docx", ".pdf", ".pptx"}
MAX_REF_FILES = 5
MAX_REF_FILE_BYTES = 5 * 1024 * 1024  # 5 MB per file
MAX_REF_TOTAL_CHARS = 30_000

MAX_SECTIONS = 10
MAX_QUESTIONS = 50
MAX_OPTIONS = 10

GEMINI_TIMEOUT = 120.0

QUESTION_TYPES = (
    "multiple_choice", "checkbox", "dropdown", "short_answer", "essay",
    "password", "date", "time", "datetime", "file_upload",
)
OPTION_TYPES = ("multiple_choice", "checkbox", "dropdown")
# quiz multiple_choice wajib tepat 1 kunci — dicek ulang saat accept.


class AiNotConfigured(Exception):
    pass


class AiFailed(Exception):
    pass


def extract_ref_text(filename: str, raw: bytes) -> str:
    """Ambil teks dari file referensi (docx/pdf/pptx)."""
    name = (filename or "").lower()
    if name.endswith(".docx"):
        from docx import Document  # type: ignore

        doc = Document(io.BytesIO(raw))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n".join(parts)
    if name.endswith(".pdf"):
        from pypdf import PdfReader  # type: ignore

        reader = PdfReader(io.BytesIO(raw))
        return "\n".join((page.extract_text() or "") for page in reader.pages)
    if name.endswith(".pptx"):
        from pptx import Presentation  # type: ignore

        prs = Presentation(io.BytesIO(raw))
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            parts.append(t)
                if shape.has_table:
                    for row in shape.table.rows:
                        parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
        return "\n".join(parts)
    raise AiFailed(f"Tipe file tidak didukung ({filename}). Pakai docx, pdf, atau pptx.")


SYSTEM_INSTRUCTION = """Kamu penyusun form/kuis berbahasa Indonesia. Jawab HANYA dengan SATU objek JSON valid, tanpa markdown, tanpa penjelasan.

Bentuk:
{"sections": [{"title": "nama section", "questions": [
  {"type": "salah satu: %s", "question_text": "teks soal",
   "is_required": true, "points": 1,
   "options": [{"option_text": "teks opsi", "is_correct": true}],
   "password_keyword": null, "answer_key": null, "allow_other": false}
]}], "settings": {"shuffle_questions": false, "shuffle_options": false, "timer_minutes": null, "require_login": false, "submission_limit": "unlimited", "show_leaderboard": false, "is_restricted": false, "show_in_history": true, "reveal_score": true, "reveal_answers": true, "display_style": "card", "scoring_mode": "auto", "theme_color": null, "thank_you_message": null, "starts_at": null, "ends_at": null}}

Aturan:
- options HANYA untuk multiple_choice/checkbox/dropdown (2-4 opsi); tipe lain: options [] dan password_keyword null.
- password_keyword HANYA untuk type password (isi kata sandinya), selain itu null.
- answer_key SELALU null — JANGAN mengarang kunci jawaban (creator mengisinya saat review; kunci salah = penilaian otomatis salah).
- allow_other HANYA true bila user EKSPILISIT meminta opsi "lainnya"/"other"/ketik-sendiri, dan HANYA untuk multiple_choice/checkbox; selain itu false.
- Quiz: multiple_choice WAJIB tepat 1 option is_correct=true; checkbox boleh >1; timer_minutes WAJIB angka 1-1440.
- Bukan quiz: timer_minutes null, is_correct semua false, show_leaderboard false, scoring_mode auto.
- submission_limit: "unlimited" atau "once" (once = wajib login, auto-coerce).
- display_style: "card" (formal) atau "quiz" (gamified satu soal per layar).
- scoring_mode: "auto" atau "manual" (hanya bermakna untuk quiz).
- theme_color: hex "#RRGGBB" atau null. thank_you_message: ringkas atau null.
- starts_at/ends_at: ISO "YYYY-MM-DDTHH:MM:SS" atau null; starts_at harus sebelum ends_at.
- question_text/option_text = teks polos, TANPA tag HTML (HTML mentah tampil sebagai teks, bukan render).
- Rumus/simbol: tulis LaTeX dengan delimiter \(...\) inline atau \[...\] display. JANGAN art Unicode (√½) dan JANGAN ejaan kata ("akar kuadrat dari").
- Kode: fence ```bahasa ... ``` (satu blok per snippet, bahasa opsional: python, javascript, java, sql, cpp, html). Kode inline: `satu backtick`.
- Link: [teks](https://...) — hanya http(s); jangan link lain.
- Maksimal 10 sections, total maksimal 30 soal. question_text ringkas.
""" % (", ".join(QUESTION_TYPES))


# prompt -> label Indonesia untuk fitur yang TIDAK didukung AI
# (banner & kategori tak pernah bisa; sisanya "diabaikan" bila AI tak menyetelnya).
IGNORED_KEYWORDS: tuple[tuple[str, tuple[str, ...]], ...] = (
    ("banner", ("banner", "cover", "gambar header", "header image")),
    ("kategori", ("kategori", "category", "kelompok form")),
)


def build_user_text(title: str, description: str | None, form_type: str, prompt: str, refs: list[tuple[str, str]]) -> str:
    parts = [
        f"Jenis: {'KUIS (ada nilai & kunci jawaban)' if form_type == 'quiz' else 'FORMULIR/pendataan (tanpa nilai)'}",
        f"Judul: {title}",
    ]
    if description:
        parts.append(f"Deskripsi: {description}")
    parts.append(f"Permintaan creator:\n{prompt}")
    for fname, text in refs:
        parts.append(f"--- Isi file referensi {fname} ---\n{text}")
    return "\n\n".join(parts)


def _gemini_models() -> list[str]:
    """Model utama + cadangan (dedupe). Kosong = tanpa fallback."""
    models = [m.strip() for m in (GEMINI_MODEL, GEMINI_FALLBACK_MODEL) if m and m.strip()]
    return list(dict.fromkeys(models)) or ["gemini-3.6-flash"]


def _repair_json_escapes(text: str) -> str:
    """Perbaiki backslash LaTeX yang tak di-escape Gemini (contoh \\( \\frac).

    JSON mode kadang mengembalikan perintah LaTeX mentah sehingga json.loads
    gagal (invalid \\escape) atau diam-diam korup (\\f jadi formfeed, \\t jadi
    tab). Aturan: perintah/delimiter LaTeX (\\frac, \\neq, \\(, ..., dikenali
    via _bfnrt_is_latex + huruf/paren/bracket) -> backslash digandakan;
    escape JSON valid (\\", \\\\, \\/, \\uXXXX, \\n/\\t/\\b/\\f/\\r yang bukan
    perintah) dibiarkan. Idempotent untuk output yang sudah benar.
    """
    out: list[str] = []
    i, n = 0, len(text)
    while i < n:
        ch = text[i]
        if ch != "\\" or i + 1 >= n:
            out.append(ch)
            i += 1
            continue
        nxt = text[i + 1]
        if nxt in '"\\/':
            out.append(text[i:i + 2])
            i += 2
            continue
        if nxt == "u":
            hexpart = text[i + 2:i + 6]
            if len(hexpart) == 4 and all(c in "0123456789abcdefABCDEF" for c in hexpart):
                out.append(text[i:i + 6])
                i += 6
            else:  # \usepackage, \underbrace, ...
                out.append("\\\\")
                i += 1
            continue
        if nxt in "bfnrt":
            if _bfnrt_is_latex(nxt, text[i + 1:]):
                out.append("\\\\")  # \neq \frac \times \right ... (LaTeX)
                i += 1
            else:
                out.append(text[i:i + 2])  # escape JSON asli (\nbaru, \ttab, ...)
                i += 2
            continue
        if nxt.isalpha() or nxt in "()[]":
            out.append("\\\\")  # perintah/delimiter LaTeX
            i += 1
            continue
        out.append(text[i:i + 2])
        i += 2
    return "".join(out)


def _parse_gemini_text(data: dict) -> dict:
    if data.get("promptFeedback", {}).get("blockReason"):
        raise AiFailed("Prompt ditolak filter keamanan AI. Coba ubah kata-katanya.")
    cands = data.get("candidates") or []
    text = (cands[0].get("content", {}).get("parts") or [{}])[0].get("text", "") if cands else ""
    try:
        parsed = json.loads(text)
    except (ValueError, TypeError):
        try:  # ponytail: 1x repair LaTeX mentah sebelum menyerah (hemat retry)
            parsed = json.loads(_repair_json_escapes(text))
            logger.info("gemini: JSON diperbaiki via escape-repair")
        except (ValueError, TypeError):
            raise AiFailed("AI gagal menyusun draf. Coba generate ulang.")
    if not isinstance(parsed, dict):
        raise AiFailed("AI gagal menyusun draf. Coba generate ulang.")
    return parsed


def call_gemini(user_text: str) -> tuple[dict, str]:
    """Panggil Gemini JSON mode. Balik (draf, model_terpakai).

    Key dikirim via header (tak muncul di URL/log). Tiap model dicoba 2x
    untuk error transient (429/5xx/network/JSON rusak); gagal semua di model
    utama -> lanjut ke fallback. 401/403 (key salah) dan 400 langsung gagal
    tanpa buang kuota coba — fallback pakai key yang sama.
    """
    if not GEMINI_API_KEY:
        raise AiNotConfigured("Fitur AI belum dikonfigurasi server.")
    payload = {
        "systemInstruction": {"parts": [{"text": SYSTEM_INSTRUCTION}]},
        "contents": [{"parts": [{"text": user_text}]}],
        "generationConfig": {"responseMimeType": "application/json", "temperature": 0.7, "maxOutputTokens": 8192},
    }
    headers = {"x-goog-api-key": GEMINI_API_KEY}
    last_err: Exception | None = None
    for model in _gemini_models():
        url = f"https://generativelanguage.googleapis.com/v1beta/models/{model}:generateContent"
        for _ in range(2):
            try:
                with httpx.Client(timeout=GEMINI_TIMEOUT) as client:
                    resp = client.post(url, headers=headers, json=payload)
            except httpx.HTTPError as e:
                logger.warning("gemini %s: network error %s", model, type(e).__name__)
                last_err = e
                continue
            if resp.status_code == 200:
                try:
                    return _parse_gemini_text(resp.json()), model
                except AiFailed as e:
                    logger.warning("gemini %s: draf tak valid (%s)", model, e)
                    last_err = e
                    continue
                except ValueError as e:
                    logger.warning("gemini %s: respons bukan JSON", model)
                    last_err = e
                    continue
            if resp.status_code in (401, 403):
                logger.error("gemini %s: key ditolak (%s)", model, resp.status_code)
                raise AiFailed("API key AI ditolak. Hubungi admin.")
            if resp.status_code == 400:
                logger.warning("gemini %s: 400 %.120s", model, resp.text)
                raise AiFailed("AI menolak permintaan. Coba ubah prompt lalu generate ulang.")
            if resp.status_code == 404:
                # ID model pensiun/diganti Google (kasus 2.x) — bukan "sibuk".
                logger.error("gemini %s: 404 model tak tersedia", model)
                last_err = AiFailed(f"Model AI {model} tidak tersedia. Hubungi admin.")
                break
            # 429 / 5xx -> coba lagi / fallback.
            logger.warning("gemini %s: sibuk (%s)", model, resp.status_code)
            last_err = AiFailed(f"AI sibuk ({resp.status_code}). Coba lagi sebentar lagi.")
    if isinstance(last_err, AiFailed):
        raise last_err
    raise AiFailed("AI tidak merespons. Periksa koneksi lalu coba lagi.")


# Konvensi rich-lite yang boleh dipakai AI di question_text/option_text:
# rumus \(...\) / \[...\] (KaTeX render di client), fence ```lang untuk kode,
# [teks](https://...) untuk link. Selain itu = teks polos.
# Pembuka fence valid: ``` + token bahasa + TERMINATOR (spasi/newline/backtick).
# Lookahead mencegah lang terpotong ("c++x..." tak jadi lang "c++") dan
# karakter asing ("evil\"...") — keduanya gagal jadi fence = teks biasa.
_FENCE_RE = re.compile(r"```([A-Za-z0-9+#_-]{0,20})(?=[ \t\n`])[ \t]*\n?(.*?)```", re.DOTALL)
_INLINE_CODE_RE = re.compile(r"`([^`\n]{1,500})`")
_LINK_RE = re.compile(r"\[([^\]\n]{1,300})\]\(([^)\s]{1,500})\)")
_URL_OK_RE = re.compile(r"^https?://[^\s<>\"]+$", re.IGNORECASE)
_LANG_OK_RE = re.compile(r"^[A-Za-z0-9+#_-]{1,20}$")

# Perintah LaTeX umum berawalan huruf escape JSON (b/f/n/r/t) — dipakai untuk
# membedakan "\neq" (LaTeX) dari "\n..." + kata (newline asli, mis. "\nbaru").
# Cocok hanya bila diikuti batas kata; sisanya dianggap escape JSON.
_LATEX_BFNRT = {
    "b": ("binom", "boldsymbol", "bigodot", "bigoplus", "bigotimes", "bigg", "bigl", "bigm", "bigr", "big", "bmod", "bot", "bowtie", "breve", "bullet", "bar", "box"),
    "f": ("frac", "fbox"),
    "n": ("notin", "nexists", "nabla", "neq", "ni"),
    "r": ("rightarrow", "rfloor", "rceil", "rangle", "right"),
    "t": ("triangle", "theta", "tilde", "times", "tau", "top", "to"),
}


def _bfnrt_is_latex(letter: str, rest: str) -> bool:
    for cmd in _LATEX_BFNRT[letter]:
        if rest.startswith(cmd) and (len(rest) == len(cmd) or not rest[len(cmd)].isalpha()):
            return True
    return False

# headroom untuk tag yang disisipkan konverter (batas kolom 5000/2000)
_RICH_HEADROOM = 500


def _inline_rich(esc: str, parts: list[str]) -> str:
    """Inline code + link di atas teks yang SUDAH di-escape. Tanpa placeholder
    bersarang: tiap temuan langsung jadi placeholder bernomor."""

    def inline_sub(m: re.Match) -> str:
        parts.append(f"<code>{m.group(1)}</code>")
        return f"\x00{len(parts) - 1}\x00"

    esc = _INLINE_CODE_RE.sub(inline_sub, esc)

    def link_sub(m: re.Match) -> str:
        url = html.unescape(m.group(2)).strip()
        if not _URL_OK_RE.match(url):
            return m.group(0)  # skema asing (javascript:/data:) = biarkan teks
        parts.append(f'<a href="{html.escape(url, quote=True)}">{m.group(1)}</a>')
        return f"\x00{len(parts) - 1}\x00"

    esc = _LINK_RE.sub(link_sub, esc)
    return esc


def _rich_lite_to_html(text: str) -> str:
    """Ubah konvensi rich-lite AI -> HTML allowlist frontend.

    Aman by construction: tiap segmen teks di-escape dulu, lalu hanya tag yang
    dibuat fungsi ini yang disisipkan (<div>/<code>/<a href http(s)>).
    HTML mentah dari AI TIDAK pernah passthrough (tampil sebagai teks).
    Delimiter LaTeX dibiarkan — KaTeX auto-render di client (pre/div kode
    dikecualikan render via ignoredTags).
    """
    parts: list[str] = []
    out: list[str] = []
    pos = 0
    for m in _FENCE_RE.finditer(text):
        out.append(_inline_rich(html.escape(text[pos:m.start()], quote=True), parts))
        lang = (m.group(1) or "").strip().lower() or "plain"
        if not _LANG_OK_RE.match(lang):
            lang = "plain"
        code = m.group(2).strip("\n")[:4000]
        parts.append(
            '<div class="ql-code-block-container">'
            f'<div class="ql-code-block" data-language="{lang}">'
            f"{html.escape(code, quote=True)}"
            "</div></div>"
        )
        out.append(f"\x00{len(parts) - 1}\x00")
        pos = m.end()
    out.append(_inline_rich(html.escape(text[pos:], quote=True), parts))
    esc = "".join(out)
    # kembalikan potongan (terdalam dulu agar placeholder bersarang aman)
    for i in range(len(parts) - 1, -1, -1):
        esc = esc.replace(f"\x00{i}\x00", parts[i])
    return esc


def _coerce_question(raw: dict) -> dict | None:
    """Bersihkan 1 soal AI -> dict valid QuestionCreate, atau None bila sampah."""
    if not isinstance(raw, dict):
        return None
    q_type = raw.get("type") if raw.get("type") in QUESTION_TYPES else None
    text = str(raw.get("question_text") or "").strip()
    if not q_type or not text:
        return None
    text = _rich_lite_to_html(text[:5000 - _RICH_HEADROOM])
    opts: list[dict] = []
    if q_type in OPTION_TYPES:
        for o in (raw.get("options") or [])[:MAX_OPTIONS]:
            if not isinstance(o, dict):
                continue
            t = str(o.get("option_text") or "").strip()
            if t:
                opts.append({"option_text": _rich_lite_to_html(t[:2000 - _RICH_HEADROOM]), "is_correct": bool(o.get("is_correct"))})
        if not opts:
            return None
    try:
        points = int(raw.get("points", 1))
    except (TypeError, ValueError):
        points = 1
    kw = str(raw.get("password_keyword") or "").strip() or None
    # answer_key: LLM dilarang mengarang kunci (lihat SYSTEM_INSTRUCTION) —
    # paksa null agar soal tak gugur validasi; creator mengisi saat review.
    # allow_other: teruskan hanya untuk MC/checkbox, selain itu False.
    allow_other = bool(raw.get("allow_other", False)) and q_type in ("multiple_choice", "checkbox")
    try:
        q = QuestionCreate(
            type=q_type,
            question_text=text,
            points=max(0, min(999, points)),
            is_required=bool(raw.get("is_required", True)),
            password_keyword=kw if q_type == "password" else None,
            answer_key=None,
            allow_other=allow_other,
            options=opts,
        )
    except Exception:
        return None
    return q.model_dump()


def sanitize_draft(raw: dict, form_type: str, prompt_text: str = "") -> dict:
    """Bersihkan output AI -> draf valid. Raise AiFailed bila tak ada soal layak."""
    sections: list[dict] = []
    total = 0
    for s in (raw.get("sections") or [])[:MAX_SECTIONS]:
        title = str((s or {}).get("title") or "").strip()[:150] or "Bagian"
        questions: list[dict] = []
        for q in ((s or {}).get("questions") or []):
            if total >= MAX_QUESTIONS:
                break
            clean = _coerce_question(q)
            if clean:
                questions.append(clean)
                total += 1
        if questions:
            sections.append({"title": title, "questions": questions})
    if not sections:
        raise AiFailed("AI tidak menghasilkan soal yang valid. Coba perjelas prompt lalu generate ulang.")

    settings = raw.get("settings") or {}
    try:
        timer = settings.get("timer_minutes")
        timer = int(timer) if timer is not None else None
        timer = timer if timer is not None and 1 <= timer <= 1440 else None
    except (TypeError, ValueError):
        timer = None
    sub = settings.get("submission_limit")
    is_quiz_type = form_type == "quiz"

    def _hex(v):
        s = str(v or "").strip()
        if len(s) == 7 and s.startswith("#"):
            try:
                int(s[1:], 16)
                return s.upper()
            except ValueError:
                return None
        return None

    def _dt(v):
        # ponytail: tanggal AI tak tentu formatnya — gagal parse = null, bukan gagal generate.
        s = str(v or "").strip()
        if not s:
            return None
        try:
            from app.schemas.form import _parse_datetime

            _parse_datetime(s)
            return s
        except Exception:
            return None

    disp = settings.get("display_style")
    disp = disp if disp in ("card", "quiz") else "card"
    scoring = settings.get("scoring_mode")
    scoring = scoring if scoring in ("auto", "manual") else "auto"
    thank = str(settings.get("thank_you_message") or "").strip()[:2000] or None
    starts = _dt(settings.get("starts_at"))
    ends = _dt(settings.get("ends_at"))
    if starts and ends:
        try:
            from app.schemas.form import _parse_datetime

            if _parse_datetime(starts) >= _parse_datetime(ends):
                starts = ends = None
        except Exception:
            starts = ends = None
    draft = {
        "sections": sections,
        "settings": {
            "shuffle_questions": bool(settings.get("shuffle_questions", False)),
            "shuffle_options": bool(settings.get("shuffle_options", False)),
            "timer_minutes": timer,
            "require_login": bool(settings.get("require_login", False)),
            "submission_limit": sub if sub in ("unlimited", "once") else "unlimited",
            "show_leaderboard": bool(settings.get("show_leaderboard", False)) and is_quiz_type,
            "is_restricted": bool(settings.get("is_restricted", False)),
            "show_in_history": bool(settings.get("show_in_history", True)),
            "reveal_score": bool(settings.get("reveal_score", True)),
            "reveal_answers": bool(settings.get("reveal_answers", True)),
            "display_style": disp,
            "scoring_mode": scoring if is_quiz_type else "auto",
            "theme_color": _hex(settings.get("theme_color")),
            "thank_you_message": thank,
            "starts_at": starts,
            "ends_at": ends,
        },
    }
    if form_type == "quiz" and timer is None:
        # Jangan diam-diam tanpa timer (publish quiz wajib timer) — creator regenerate.
        raise AiFailed("AI tidak menyertakan timer untuk kuis. Coba generate ulang.")
    draft["ignored"] = detect_ignored(prompt_text, draft["settings"])
    return draft


def detect_ignored(prompt_text: str, settings: dict) -> list[str]:
    """Minta user menyebut fitur X tapi AI tak menyetelnya -> label untuk warning box.

    Murni heuristik substring (ID+EN), tanpa panggilan AI tambahan (hemat kuota).
    """
    p = (prompt_text or "").lower()
    out: list[str] = []

    def has(*words: str) -> bool:
        return any(w in p for w in words)

    for label, words in IGNORED_KEYWORDS:
        if has(*words):
            out.append(label)
    s = settings or {}
    if has("leaderboard", "peringkat", "papan skor", "ranking") and not s.get("show_leaderboard"):
        out.append("leaderboard")
    if has("warna", "theme", "colour", "color", "ungu", "biru", "merah", "hijau") and not s.get("theme_color"):
        out.append("warna tema")
    if has("jadwal", "dibuka", "ditutup", "berakhir", "deadline", "batas waktu", "tanggal mulai", "tanggal selesai") and not (s.get("starts_at") or s.get("ends_at")):
        out.append("jadwal")
    if has("terima kasih", "thank you", "pesan penutup", "closing message") and not s.get("thank_you_message"):
        out.append("pesan terima kasih")
    if has("tampilan kuis", "gamified", "satu soal per layar", "mode kuis") and s.get("display_style") == "card":
        out.append("tampilan")
    if has("skor manual", "bobot nilai", "penilaian manual", "manual scoring") and s.get("scoring_mode") == "auto":
        out.append("mode penilaian")
    if has("sembunyikan skor", "sembunyikan nilai", "tanpa skor") and s.get("reveal_score"):
        out.append("tampil skor")
    if has("sembunyikan kunci", "sembunyikan jawaban", "tanpa pembahasan") and s.get("reveal_answers"):
        out.append("tampil jawaban")
    if has("sembunyikan dari riwayat", "tanpa riwayat") and s.get("show_in_history"):
        out.append("riwayat")
    if has("terbatas", "restricted", "hanya undangan") and not s.get("is_restricted"):
        out.append("mode terbatas")
    # dedupe, jaga urutan
    return list(dict.fromkeys(out))
